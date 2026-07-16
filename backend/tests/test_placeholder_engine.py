from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.services.pptx.placeholder_engine import (
    find_special_placeholders,
    insert_picture,
    replace_text_placeholders,
)
from tests import fixtures


def _all_text(slide) -> str:
    out = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            out.append(shape.text_frame.text)
    return "\n".join(out)


def _find_run(slide, contains: str):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if contains in run.text:
                    return run
    return None


def test_replacement_preserves_formatting_across_split_runs(tmp_path: Path) -> None:
    path = fixtures.make_template_pptx(tmp_path / "tmpl.pptx")
    prs = Presentation(str(path))
    slide = prs.slides[0]

    mapping = {
        "Title": "Acme Manufacturing Transformation",
        "Customer": "Acme Corp",
        "Industry": "Manufacturing",
        "Region": "EMEA",
        "Challenge": "Fragmented data",
        "Solution": "SAP S/4HANA rollout",
        "Benefits": "• 30% savings",
    }
    replace_text_placeholders(slide.shapes, mapping)

    text = _all_text(slide)
    # All text-field tokens are replaced; only the whole-shape {{Image}} special
    # remains (handled separately by the deck assembler).
    for tok in ("{{Title}}", "{{Customer}}", "{{Industry}}", "{{Region}}", "{{Challenge}}"):
        assert tok not in text
    assert text.count("{{") == 1 and "{{Image}}" in text
    assert "Acme Manufacturing Transformation" in text
    assert "Client: Acme Corp" in text
    assert "Industry: Manufacturing | Region: EMEA" in text

    # The split {{Title}} token: replacement lands in the anchor run and keeps format.
    run = _find_run(slide, "Acme Manufacturing Transformation")
    assert run is not None
    assert run.font.bold is True
    assert run.font.size is not None and run.font.size.pt == 28
    assert run.font.color.rgb is not None and str(run.font.color.rgb) == "FF0000"

    # Customer run keeps its font name + size.
    cust_run = _find_run(slide, "Acme Corp")
    assert cust_run is not None
    assert cust_run.font.name == "Georgia"
    assert cust_run.font.size.pt == 18


def test_image_placeholder_inserts_picture_and_preserves_geometry(tmp_path: Path) -> None:
    path = fixtures.make_template_pptx(tmp_path / "tmpl.pptx")
    prs = Presentation(str(path))
    slide = prs.slides[0]

    specials = find_special_placeholders(slide)
    assert len(specials) == 1
    ph = specials[0]
    assert ph.kind == "image"
    box = (ph.left, ph.top, ph.width, ph.height)

    img_path = tmp_path / "logo.png"
    img_path.write_bytes(fixtures._png_bytes())
    insert_picture(slide, ph, str(img_path))

    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1
    pic = pics[0]
    # Picture fits within the original placeholder box.
    assert pic.left >= box[0]
    assert pic.top >= box[1]
    assert pic.width <= box[2] + 1
    # The {{Image}} text placeholder shape was removed.
    assert "{{Image}}" not in _all_text(slide)


def test_reopens_as_valid_pptx(tmp_path: Path) -> None:
    path = fixtures.make_template_pptx(tmp_path / "tmpl.pptx")
    prs = Presentation(str(path))
    replace_text_placeholders(prs.slides[0].shapes, {"Title": "X", "Customer": "Y"})
    out = tmp_path / "out.pptx"
    prs.save(str(out))
    # Re-open to confirm structural validity.
    reopened = Presentation(str(out))
    assert len(reopened.slides) == 1
