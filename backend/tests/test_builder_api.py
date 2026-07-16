from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from tests import fixtures


def _upload_doc(client, path: Path):
    with path.open("rb") as fh:
        return client.post("/api/documents", files=[("files", (path.name, fh.read()))])


def _upload_template(client, path: Path, name: str = "Single Case Study"):
    with path.open("rb") as fh:
        return client.post(
            "/api/templates",
            files=[("file", (path.name, fh.read()))],
            data={"name": name, "category": "Manufacturing"},
        )


def test_template_upload_discovers_placeholders(client, tmp_path: Path) -> None:
    p = fixtures.make_template_pptx(tmp_path / "tmpl.pptx")
    resp = _upload_template(client, p)
    assert resp.status_code == 200, resp.text
    t = resp.json()
    names = {ph["name"] for ph in t["placeholders"]}
    assert "Title" in names
    assert "Customer" in names
    assert "Image" in names
    kinds = {ph["name"]: ph["kind"] for ph in t["placeholders"]}
    assert kinds["Image"] == "image"
    assert kinds["Title"] == "text"
    assert t["slide_count"] == 1


def test_build_presentation_from_selection(client, tmp_path: Path) -> None:
    # Seed a manufacturing case study.
    doc_path = fixtures.make_txt(
        tmp_path / "acme.txt",
        text="Acme Manufacturing modernized its ERP with SAP S/4HANA, cutting "
        "inventory costs 30% across EMEA in 6 months.",
    )
    doc = _upload_doc(client, doc_path).json()[0]

    # Find its case study id via search.
    results = client.post("/api/search", json={"query": "Acme SAP manufacturing"}).json()["items"]
    assert results
    cs_id = results[0]["case_study_id"]

    tmpl = _upload_template(client, fixtures.make_template_pptx(tmp_path / "tmpl.pptx")).json()

    resp = client.post(
        "/api/presentations",
        json={
            "name": "Manufacturing Deck",
            "template_id": tmpl["id"],
            "case_study_ids": [cs_id],
            "options": {"include_agenda": True, "include_thank_you": True},
        },
    )
    assert resp.status_code == 200, resp.text
    pres = resp.json()
    assert pres["status"] == "ready"
    assert pres["output_pptx_path"]

    # Validate the produced PPTX: tokens filled, title present, > 1 slide (title/agenda/thankyou).
    out = Path(pres["output_pptx_path"])
    assert out.exists()
    prs = Presentation(str(out))
    assert len(prs.slides) >= 3
    all_text = "\n".join(
        s.text_frame.text
        for slide in prs.slides
        for s in slide.shapes
        if s.has_text_frame
    )
    assert "{{" not in all_text
    assert "Acme" in all_text
    assert "Thank You" in all_text

    _ = doc  # keep reference


def test_build_from_prompt(client, tmp_path: Path) -> None:
    doc_path = fixtures.make_txt(
        tmp_path / "g.txt",
        text="Globex Healthcare deployed an AI diagnostics platform in APAC, "
        "reducing patient wait times 40% in 9 months.",
    )
    _upload_doc(client, doc_path)
    tmpl = _upload_template(client, fixtures.make_template_pptx(tmp_path / "tmpl.pptx")).json()

    resp = client.post(
        "/api/presentations",
        json={
            "name": "Healthcare Deck",
            "template_id": tmpl["id"],
            "prompt": "Create a healthcare presentation for AI implementations, top 3",
            "options": {"max_case_studies": 3},
        },
    )
    assert resp.status_code == 200, resp.text
    pres = resp.json()
    assert pres["status"] == "ready"
    assert pres["intent"]["industries"] == ["Healthcare"]
    assert len(pres["selected_case_study_ids"]) >= 1
