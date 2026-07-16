"""Helpers that programmatically generate small sample documents for tests.

Keeping fixtures generated (rather than committing binaries) keeps the repo small
and makes the expected content explicit.
"""

from __future__ import annotations

import io
from pathlib import Path

SAMPLE_TEXT = (
    "Acme Manufacturing modernized its ERP with SAP S/4HANA. "
    "The challenge was fragmented supply chain data across EMEA. "
    "The solution delivered a 30% reduction in inventory costs over 6 months."
)


def _png_bytes(color: tuple[int, int, int] = (200, 100, 50)) -> bytes:
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", (64, 48), color).save(buf, format="PNG")
    return buf.getvalue()


def make_txt(path: Path, text: str = SAMPLE_TEXT) -> Path:
    path.write_text(text, encoding="utf-8")
    return path


def make_docx(path: Path, text: str = SAMPLE_TEXT) -> Path:
    import docx

    doc = docx.Document()
    doc.add_heading("Acme Manufacturing Case Study", level=1)
    doc.add_paragraph(text)
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Inventory reduction"
    table.cell(1, 1).text = "30%"
    doc.save(str(path))
    return path


def make_pptx(path: Path, text: str = SAMPLE_TEXT) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Acme Manufacturing Case Study"
    slide.placeholders[1].text = text

    # A slide with a table + image.
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    rows, cols = 2, 2
    tbl = slide2.shapes.add_table(
        rows, cols, Inches(1), Inches(1), Inches(4), Inches(1.2)
    ).table
    tbl.cell(0, 0).text = "Benefit"
    tbl.cell(0, 1).text = "Impact"
    tbl.cell(1, 0).text = "Cost savings"
    tbl.cell(1, 1).text = "30%"

    img_path = path.parent / "_tmp_img.png"
    img_path.write_bytes(_png_bytes())
    slide2.shapes.add_picture(str(img_path), Inches(6), Inches(1), Inches(2), Inches(1.5))

    prs.save(str(path))
    img_path.unlink(missing_ok=True)
    return path


def make_template_pptx(path: Path) -> Path:
    """Build a template with mixed-formatting placeholders, including a token that is
    deliberately split across two runs (as PowerPoint often does) and an image
    placeholder shape."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title textbox with a run-split {{Title}} token, bold + red + size 28.
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
    p = title_box.text_frame.paragraphs[0]
    r1 = p.add_run()
    r1.text = "{{Ti"
    r1.font.bold = True
    r1.font.size = Pt(28)
    r1.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    r2 = p.add_run()
    r2.text = "tle}}"
    r2.font.bold = True
    r2.font.size = Pt(28)
    r2.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)

    # Customer textbox, single run, size 18, named font.
    cust = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.6))
    cp = cust.text_frame.paragraphs[0]
    cr = cp.add_run()
    cr.text = "Client: {{Customer}}"
    cr.font.size = Pt(18)
    cr.font.name = "Georgia"

    # Two tokens in one paragraph.
    meta = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(0.6))
    mp = meta.text_frame.paragraphs[0]
    mr = mp.add_run()
    mr.text = "Industry: {{Industry}} | Region: {{Region}}"
    mr.font.size = Pt(14)

    # Body with challenge/solution.
    body = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9), Inches(2))
    for label in ("Challenge: {{Challenge}}", "Solution: {{Solution}}", "Benefits:\n{{Benefits}}"):
        para = body.text_frame.add_paragraph()
        rr = para.add_run()
        rr.text = label
        rr.font.size = Pt(12)

    # Image placeholder (whole-shape token).
    img = slide.shapes.add_textbox(Inches(6.5), Inches(1.2), Inches(2.5), Inches(1.8))
    img.text_frame.paragraphs[0].add_run().text = "{{Image}}"

    prs.save(str(path))
    return path


def make_pdf(path: Path, text: str = SAMPLE_TEXT) -> Path:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    c = canvas.Canvas(str(path), pagesize=letter)
    c.setFont("Helvetica-Bold", 16)
    c.drawString(72, 720, "Acme Manufacturing Case Study")
    c.setFont("Helvetica", 11)
    y = 690
    for line in _wrap(text, 80):
        c.drawString(72, y, line)
        y -= 16
    c.showPage()
    c.save()
    return path


def _wrap(text: str, width: int) -> list[str]:
    words = text.split()
    lines: list[str] = []
    cur = ""
    for w in words:
        if len(cur) + len(w) + 1 > width:
            lines.append(cur)
            cur = w
        else:
            cur = f"{cur} {w}".strip()
    if cur:
        lines.append(cur)
    return lines
