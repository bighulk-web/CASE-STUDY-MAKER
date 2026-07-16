"""Discover ``{{placeholders}}`` in a PowerPoint template."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.services.pptx.placeholder_engine import SPECIAL_KINDS, TOKEN_RE


def _iter_text(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_text(shape.shapes)
            continue
        if shape.has_text_frame:
            yield shape.text_frame.text
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text


def discover_placeholders(path: str | Path) -> tuple[list[dict[str, object]], int]:
    """Return (placeholders, slide_count).

    Each placeholder: ``{name, kind, arg, slides:[int]}`` where ``kind`` is ``text``
    or one of the special kinds (image/logo/icon/table/chart).
    """
    prs = Presentation(str(path))
    found: dict[str, dict[str, object]] = {}
    slide_count = 0

    for idx, slide in enumerate(prs.slides):
        slide_count += 1
        for text in _iter_text(slide.shapes):
            for m in TOKEN_RE.finditer(text or ""):
                token = m.group(1).strip()
                base = token.split(":", 1)[0].strip()
                arg = token.split(":", 1)[1].strip() if ":" in token else ""
                kind = base.lower() if base.lower() in SPECIAL_KINDS else "text"
                entry = found.setdefault(
                    token, {"name": token, "kind": kind, "arg": arg, "slides": []}
                )
                slides = entry["slides"]
                assert isinstance(slides, list)
                if idx not in slides:
                    slides.append(idx)

    return list(found.values()), slide_count
