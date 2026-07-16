"""PowerPoint (.pptx) extractor using python-pptx."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .base import ExtractedChart, ExtractedImage, ExtractedTable, ExtractionResult


def _shape_text(shape) -> str:
    if shape.has_text_frame:
        return "\n".join(p.text for p in shape.text_frame.paragraphs)
    return ""


def _extract_table(shape) -> ExtractedTable:
    tbl = shape.table
    rows = [[cell.text for cell in row.cells] for row in tbl.rows]
    return ExtractedTable(rows=rows)


def _extract_chart(shape) -> ExtractedChart:
    chart = shape.chart
    out = ExtractedChart()
    try:
        out.title = chart.chart_title.text_frame.text if chart.has_title else ""
    except Exception:
        out.title = ""
    try:
        plot = chart.plots[0]
        out.categories = [str(c) for c in plot.categories]
    except Exception:
        pass
    try:
        for series in chart.series:
            out.series[str(series.name)] = [float(v) if v is not None else 0.0 for v in series.values]
    except Exception:
        pass
    return out


def _walk(shapes, result: ExtractionResult, texts: list[str]) -> None:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _walk(shape.shapes, result, texts)
            continue
        if shape.has_table:
            result.tables.append(_extract_table(shape))
            continue
        if getattr(shape, "has_chart", False) and shape.has_chart:
            result.charts.append(_extract_chart(shape))
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                img = shape.image
                result.images.append(ExtractedImage(data=img.blob, ext=img.ext or "png"))
            except Exception:
                pass
            continue
        t = _shape_text(shape)
        if t.strip():
            texts.append(t)


def extract_pptx(path: Path) -> ExtractionResult:
    prs = Presentation(str(path))
    result = ExtractionResult()
    texts: list[str] = []
    slide_count = 0
    for slide in prs.slides:
        slide_count += 1
        _walk(slide.shapes, result, texts)
        # Slide notes add valuable context.
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                texts.append(notes)
    result.text = "\n\n".join(texts)
    result.page_count = slide_count
    return result
